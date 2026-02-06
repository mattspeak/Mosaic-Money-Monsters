import os
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import tempfile

# Configuration
HTML_FILE = '/Users/mattspeak/Library/Mobile Documents/iCloud~md~obsidian/Documents/ObsidianMatt/WORK/Mosaic/Money Monsters User Testing Feb-26/Research/Desk_research/innovators_report.html'
OUTPUT_FILE = '/Users/mattspeak/Library/Mobile Documents/iCloud~md~obsidian/Documents/ObsidianMatt/WORK/Mosaic/Money Monsters User Testing Feb-26/Research/Desk_research/Innovators_Deck.pptx'

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to Widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def download_image(url, base_url=None):
    try:
        # Resolve relative URLs
        if url.startswith('/'):
            if base_url:
                url = base_url.rstrip('/') + url
            else:
                print(f"Skipping relative URL without base: {url}")
                return None
            
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            image_stream = io.BytesIO(response.content)
            # Check format and convert if necessary (e.g. WebP to PNG)
            try:
                from PIL import Image
                img = Image.open(image_stream)
                if img.format not in ['PNG', 'JPEG', 'GIF', 'BMP', 'TIFF']:
                    # Convert to PNG
                    new_stream = io.BytesIO()
                    img.convert('RGB').save(new_stream, format='PNG')
                    new_stream.seek(0)
                    return new_stream
                else:
                    image_stream.seek(0)
                    return image_stream
            except Exception as e:
                print(f"Error processing image format for {url}: {e}")
                # Return original stream and hope for the best if PIL fails, though pptx will likely reject it
                image_stream.seek(0)
                return image_stream

    except Exception as e:
        print(f"Error downloading image {url}: {e}")
    return None

def add_company_slide(prs, company_data):
    # Blank slide layout
    blank_slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(blank_slide_layout)

    # --- 1. Header Section ---
    # Company Name (Left formatted)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(10)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = company_data['name']
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(44, 62, 80) # Dark Blue/Grey

    # Logo (Right formatted)
    if company_data['logo_url']:
        logo_stream = download_image(company_data['logo_url'], base_url=company_data['url'])
        if logo_stream:
            try:
                # Place top right
                logo_size = Inches(0.8)
                slide.shapes.add_picture(logo_stream, Inches(12), Inches(0.5), width=logo_size)
            except Exception as e:
                print(f"Could not add logo: {e}")

    # Separator Line
    shape = slide.shapes.add_connector(
        1, Inches(0.5), Inches(1.6), Inches(12.83), Inches(1.6)
    )
    shape.line.color.rgb = RGBColor(200, 200, 200)
    shape.line.width = Pt(2)

    # --- 2. Main Content Columns ---
    
    # URL underneath title
    txBoxUrl = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(0.4))
    tfUrl = txBoxUrl.text_frame
    pUrl = tfUrl.paragraphs[0]
    pUrl.text = company_data['url'].strip()
    pUrl.font.size = Pt(12)
    pUrl.font.color.rgb = RGBColor(52, 152, 219) # Blue

    # Left Column: Intro Copy + Main Image
    left_col_x = Inches(0.5)
    left_col_y = Inches(2.0)
    left_col_w = Inches(5.0)

    # Intro Text
    txBoxIntro = slide.shapes.add_textbox(left_col_x, left_col_y, left_col_w, Inches(2))
    tfIntro = txBoxIntro.text_frame
    tfIntro.word_wrap = True
    pIntro = tfIntro.paragraphs[0]
    pIntro.text = company_data['intro']
    pIntro.font.size = Pt(18)
    pIntro.font.name = 'Arial'
    
    image_y = Inches(4.5)
    if company_data['image_url']:
        img_stream = download_image(company_data['image_url'], base_url=company_data['url'])
        if img_stream:
            try:
                # Add picture, constrain width to column width
                pic = slide.shapes.add_picture(img_stream, left_col_x, image_y, width=left_col_w)
            except Exception as e:
                print(f"Could not add main image: {e}")

    # Right Column: Supporting Principles / Details
    right_col_x = Inches(6.0)
    right_col_y = Inches(2.0)
    right_col_w = Inches(6.8)

    txBoxDetails = slide.shapes.add_textbox(right_col_x, right_col_y, right_col_w, Inches(5))
    tfDetails = txBoxDetails.text_frame
    tfDetails.word_wrap = True
    
    # Add bullets
    if company_data.get('bullets'):
        for bullet_text in company_data['bullets']:
            p = tfDetails.add_paragraph()
            p.text = bullet_text
            p.font.size = Pt(14)
            p.font.name = 'Arial'
            p.level = 0 # Top level bullet
            # Fix manual bullets if they exist in text
            if p.text.startswith("- ") or p.text.startswith("• "):
                 p.text = p.text[2:]
            
            # Simple space between bullets
            p.space_after = Pt(6)
    else:
        # Fallback if no bullets found
        p = tfDetails.add_paragraph()
        p.text = company_data.get('details', '')
        p.font.size = Pt(14)

def clean_text(text):
    if not text: return ""
    import re
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def parse_html_report(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')

    companies = []
    cards = soup.find_all('div', class_='company-card')

    for card in cards:
        # Title
        h2 = card.find('h2')
        name = h2.get_text(strip=True) if h2 else "Unknown Organization"

        # URL
        a_tag = card.find('a', class_='company-url')
        url = a_tag['href'] if a_tag else ""

        # Logo
        logo_img = card.find('img', class_='company-logo')
        logo_url = logo_img['src'] if logo_img else None

        # Main Image
        main_img = card.find('img', class_='main-image')
        image_url = main_img['src'] if main_img else None

        # Content Parsing
        summary_div = card.find('div', class_='summary')
        intro_text = ""
        bullets = []

        if summary_div:
            # Plan:
            # 1. Convert div to text, but keeping some delimiters to help splitting.
            # 2. Or iterate children.
            
            # Let's use a keyword-based split on the full text first to find the "Break Point"
            # Known headers: "Key Offerings", "Key Features", "Insights on", "Principles of", "Summary:"
            # Actually, "Summary:" usually starts the intro.
            # The "Details" start with "Key Offerings", "Key Features", "Insights...", "Principles..."
            
            # Let's try to identify the *node* that starts the details section.
            
            details_keywords = ["key offerings", "key features", "insights", "principles", "success principles"]
            
            intro_nodes = []
            details_nodes = []
            mode = 'intro'
            
            for child in summary_div.children:
                text_content = child.get_text(strip=True).lower()
                
                # Check for header switch
                is_header = False
                if child.name in ['h4', 'h3', 'strong', 'b']:
                    if any(k in text_content for k in details_keywords):
                        is_header = True
                
                # Also check plain text if it looks like a header (short, specific words)
                if not is_header and len(text_content) < 50:
                     if any(k in text_content for k in details_keywords):
                         is_header = True
                
                if is_header:
                    mode = 'details'
                
                if mode == 'intro':
                    intro_nodes.append(child)
                else:
                    details_nodes.append(child)
            
            # Process Intro
            intro_accum = []
            for node in intro_nodes:
                if isinstance(node, str):
                    t = node.strip()
                    if t: intro_accum.append(t)
                else:
                    t = node.get_text(strip=True)
                    if t: intro_accum.append(t)
            
            intro_full = " ".join(intro_accum)
            # Cleanup "Summary:" prefix if present
            if intro_full.lower().startswith("summary"):
                intro_full = intro_full.split(":", 1)[-1].strip()
            # Cleanup bold tags in text (e.g. Summary: Title)
            intro_full = clean_text(intro_full)
            
            # Process Details into Bullets
            # We want to iterate through details_nodes and extract list items.
            # Note: The HTML is often just text with <br>.
            
            raw_details_text = ""
            for node in details_nodes:
                # Get text with some separator preservation?
                # If we just get_text, we lose <br> boundaries which are critical here.
                if hasattr(node, 'get_text'):
                    # Custom stringify to keep newlines for <br>
                    s = str(node)
                    s = s.replace('<br>', '\n').replace('<br/>', '\n')
                    # Remove tags
                    import re
                    s = re.sub(r'<[^>]+>', '', s)
                    raw_details_text += s + "\n"
                elif isinstance(node, str):
                    raw_details_text += node + "\n"
            
            # Now split by lines and look for bullets
            lines = raw_details_text.split('\n')
            for line in lines:
                line = line.strip()
                if not line: continue
                
                # Skip the header lines themselves if they are purely headers
                if any(line.lower().startswith(k) for k in details_keywords) and len(line) < 60:
                    continue
                    
                # Identify bullet pattern
                # 1. Start with digit + dot/paren (1. , 1) )
                # 2. Start with - or •
                import re
                
                # Remove leading numbering/bullet chars
                cleaned = re.sub(r'^(\d+[\.\)]|\-|\•)\s*', '', line)
                
                # Heuristic: If it's short, or starts with known pattern, it's a bullet.
                # If it's a long paragraph, maybe split it?
                # User wants "Sharp bullets".
                
                # Check for "Key: Value" pattern (e.g. "Money Wise Game: An interactive...")
                # We can bold the key if we want, but python-pptx simple text run doesn't mix easily in one paragraph efficiently without more code.
                # Let's just keep the text.
                
                bullets.append(cleaned)

        companies.append({
            'name': name,
            'url': url,
            'logo_url': logo_url,
            'image_url': image_url,
            'intro': intro_text if intro_text else "No introduction available.",
            'bullets': bullets,
            'details': "" # deprecated
        })

    return companies

def main():
    print("Parsing HTML report...")
    if not os.path.exists(HTML_FILE):
        print(f"Error: HTML file not found at {HTML_FILE}")
        return

    companies = parse_html_report(HTML_FILE)
    print(f"Found {len(companies)} companies.")

    print("Creating presentation...")
    prs = create_presentation()

    for i, company in enumerate(companies):
        print(f"Processing slide {i+1}: {company['name']}")
        add_company_slide(prs, company)

    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to: {OUTPUT_FILE}")

if __name__ == "__main__":
    main()
