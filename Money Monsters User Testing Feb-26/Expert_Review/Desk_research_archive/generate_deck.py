"""
Generate PowerPoint deck: Financial Literacy Landscape Analysis
For Money Monsters User Testing Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(44, 62, 80)
ACCENT_BLUE = RGBColor(52, 152, 219)
LIGHT_GRAY = RGBColor(236, 240, 241)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(230, 126, 34)
GREEN = RGBColor(39, 174, 96)
PURPLE = RGBColor(142, 68, 173)
RED = RGBColor(192, 57, 43)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_header(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), Inches(13.333), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, notes=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(14)

    return slide

def add_framework_slide(prs, title, directions):
    """Add a slide showing the four directions framework"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Four boxes for directions
    colors = [ORANGE, ACCENT_BLUE, GREEN, PURPLE]
    positions = [
        (0.5, 1.6),    # Top left
        (6.916, 1.6),  # Top right
        (0.5, 4.4),    # Bottom left
        (6.916, 4.4),  # Bottom right
    ]

    for i, (direction, color) in enumerate(zip(directions, colors)):
        x, y = positions[i]

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.917), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Direction label
        label_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(5.5), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = direction['label']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Direction name
        name_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.6), Inches(5.5), Inches(0.6))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = direction['name']
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Tagline
        tag_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.2), Inches(5.5), Inches(0.5))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = f'"{direction["tagline"]}"'
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = WHITE

        # Best for
        best_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.8), Inches(5.5), Inches(0.5))
        tf = best_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Best for: {direction['best_for']}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

    return slide

def add_direction_detail_slide(prs, direction_data, color):
    """Add detailed slide for a specific direction"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar with direction color
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{direction_data['label']}: {direction_data['name']}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(6), Inches(0.6))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = f'"{direction_data["tagline"]}"'
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = color

    # Key mechanics (left column)
    mech_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.8), Inches(4))
    tf = mech_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Mechanics"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    for mechanic in direction_data['mechanics']:
        p = tf.add_paragraph()
        p.text = "• " + mechanic
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)

    # Examples (right column)
    ex_box = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.8), Inches(4))
    tf = ex_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Examples"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    for example in direction_data['examples']:
        p = tf.add_paragraph()
        p.text = "• " + example
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)

    # Best for footer
    best_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.5))
    tf = best_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Best audience fit: {direction_data['best_for']}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color

    return slide

def add_modality_headers_slide(prs):
    """Add slide showing the modality headers for sorting apps"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Four Approaches: How to Sort Solutions"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    modalities = [
        ("TEACH: App-based", "Gamified apps with progressive lessons, quizzes, badges", "TEACH"),
        ("TEACH: Institutional", "School/employer delivered, curriculum-aligned", "TEACH"),
        ("TEACH: Simulation", "Games and simulations, low-stakes practice", "TEACH"),
        ("PRACTICE", "Real money management: banking, investing, saving", "PRACTICE"),
        ("CONNECT", "Social learning: peers, finfluencers, communities", "CONNECT"),
        ("GUIDE", "Personalized support: AI coaches, adaptive learning", "GUIDE"),
    ]

    y_start = 1.5
    for i, (name, desc, direction) in enumerate(modalities):
        y = y_start + (i * 0.95)

        # Number/direction indicator
        num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.8), Inches(0.7))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = ACCENT_BLUE
        num_box.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.15), Inches(0.8), Inches(0.5))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = direction
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Name
        name_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(4.5), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Description
        desc_box = slide.shapes.add_textbox(Inches(6.2), Inches(y), Inches(6.5), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE

    return slide

def add_apps_by_modality_slide(prs, modality_name, apps, color):
    """Add slide listing apps for a specific modality"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = modality_name
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Apps in columns
    col_width = 6
    apps_per_col = (len(apps) + 1) // 2

    for col in range(2):
        x = 0.5 + (col * 6.5)
        start_idx = col * apps_per_col
        end_idx = min(start_idx + apps_per_col, len(apps))

        y = 1.5
        for app in apps[start_idx:end_idx]:
            # App box
            app_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6), Inches(1.1))
            app_box.fill.solid()
            app_box.fill.fore_color.rgb = LIGHT_GRAY
            app_box.line.fill.background()

            # App name
            name_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(5.6), Inches(0.4))
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = app['name']
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE

            # Description
            desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.5), Inches(5.6), Inches(0.6))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = app['desc']
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_BLUE

            y += 1.2

    return slide

def add_audience_fit_slide(prs):
    """Add audience fit analysis slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Audience Fit: Which Approach for Which Age?"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Headers row
    headers = ["Approach", "Young Teens\n(13-15)", "Older Teens\n(16-18)", "Young Adults\n(18-25)"]
    col_widths = [3.5, 3, 3, 3]
    x_start = 0.5

    for i, (header, width) in enumerate(zip(headers, col_widths)):
        x = x_start + sum(col_widths[:i])
        header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(width - 0.1), Inches(0.8))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = ACCENT_BLUE
        header_box.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x), Inches(1.55), Inches(width - 0.1), Inches(0.7))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    rows = [
        ("TEACH", "Strong", "Medium", "Weak"),
        ("PRACTICE", "Weak", "Medium", "Strong"),
        ("CONNECT", "Strong", "Strong", "Medium"),
        ("GUIDE", "Medium", "Strong", "Strong"),
    ]

    colors_map = {"Strong": GREEN, "Medium": ORANGE, "Weak": RED}
    row_colors = [ORANGE, PURPLE, GREEN, ACCENT_BLUE]

    y = 2.4
    for row_idx, (direction, *fits) in enumerate(rows):
        # Direction column
        x = x_start
        dir_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_widths[0] - 0.1), Inches(0.9))
        dir_box.fill.solid()
        dir_box.fill.fore_color.rgb = row_colors[row_idx]
        dir_box.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.25), Inches(col_widths[0] - 0.2), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = direction
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Fit columns
        for col_idx, fit in enumerate(fits):
            x = x_start + sum(col_widths[:col_idx + 1])
            fit_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_widths[col_idx + 1] - 0.1), Inches(0.9))
            fit_box.fill.solid()
            fit_box.fill.fore_color.rgb = LIGHT_GRAY
            fit_box.line.fill.background()

            # Colored indicator
            indicator = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(y + 0.3), Inches(0.3), Inches(0.3))
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = colors_map[fit]
            indicator.line.fill.background()

            text_box = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y + 0.25), Inches(col_widths[col_idx + 1] - 0.9), Inches(0.5))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = fit
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_BLUE

        y += 1.0

    # Key insight
    insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12), Inches(0.8))
    tf = insight_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key insight: TEACH (gamified) works for young teens but becomes a turn-off for 18-25s who need PRACTICE or GUIDE"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = DARK_BLUE

    return slide

def add_so_what_slide(prs):
    """Add the 'So What' implications slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "So What: Implications for Money Monsters"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Current position
    current_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(2))
    current_box.fill.solid()
    current_box.fill.fore_color.rgb = ORANGE
    current_box.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(1.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Where Money Monsters is today"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "TEACH (app-based, gamified)"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "• Progressive curriculum with levels"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "• Monster-taming gamification"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "• Strong for young teen acquisition"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Challenge
    challenge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.5), Inches(6), Inches(2))
    challenge_box.fill.solid()
    challenge_box.fill.fore_color.rgb = RED
    challenge_box.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(7.033), Inches(1.6), Inches(5.6), Inches(1.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The challenge"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "TEACH alone is insufficient for:"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "• Older users (16-25) who reject gamification"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "• Real behaviour change"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "• Long-term retention"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Evolution path
    path_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(12.333), Inches(2.8))
    path_box.fill.solid()
    path_box.fill.fore_color.rgb = GREEN
    path_box.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(12), Inches(2.6))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Recommended evolution path"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p = tf.add_paragraph()
    p.text = "Phase 1: Strengthen TEACH with content partnerships (CONNECT elements)"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Phase 2: Add GUIDE layer (AI coach) — the unifying element that works across all ages"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Phase 3: Add PRACTICE (banking partnerships) OR expand CONNECT (social features)"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

    return slide

def add_engagement_principles_slide(prs):
    """Add slide showing the four engagement principles"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Four Engagement Principles"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Four principle boxes
    principles = [
        {
            'num': '1',
            'name': 'PROPOSITION',
            'question': '"What\'s in it for me?"',
            'desc': 'Clear, tangible outcome the learner will gain',
            'color': ORANGE
        },
        {
            'num': '2',
            'name': 'UTILITY',
            'question': '"Does this interaction have a point?"',
            'desc': 'Every interactive element drives toward a functional outcome',
            'color': ACCENT_BLUE
        },
        {
            'num': '3',
            'name': 'CONTEXTUAL FRAMING',
            'question': '"Does this connect to my world?"',
            'desc': 'Same concept, presented through contexts they recognize',
            'color': GREEN
        },
        {
            'num': '4',
            'name': 'CAPABILITY',
            'question': '"Am I being talked down to?"',
            'desc': 'Don\'t underestimate — sophisticated concepts, made relational',
            'color': PURPLE
        }
    ]

    positions = [
        (0.5, 1.5),    # Top left
        (6.916, 1.5),  # Top right
        (0.5, 4.3),    # Bottom left
        (6.916, 4.3),  # Bottom right
    ]

    for i, (principle, (x, y)) in enumerate(zip(principles, positions)):
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.917), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = principle['color']
        box.line.fill.background()

        # Number
        num_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = principle['num']
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Name
        name_box = slide.shapes.add_textbox(Inches(x + 0.6), Inches(y + 0.15), Inches(5), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = principle['name']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Question
        q_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(5.5), Inches(0.5))
        tf = q_box.text_frame
        p = tf.paragraphs[0]
        p.text = principle['question']
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.3), Inches(5.5), Inches(1))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = principle['desc']
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

    return slide

def add_action_1_slide(prs):
    """Action 1: Lead with Outcomes (PROPOSITION)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar - Orange for PROPOSITION
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Action 1: Lead with Outcomes"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Principle tag
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Based on: PROPOSITION — \"What's in it for me?\""
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = ORANGE

    # The action
    action_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.1), Inches(12.333), Inches(1.2))
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = LIGHT_GRAY
    action_box.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(1))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Every lesson/level should open with a clear, tangible benefit"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p = tf.add_paragraph()
    p.text = "Not just \"Learn about X\" — but \"After this, you'll be able to Y\""
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE

    # Template section
    template_label = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.4))
    tf = template_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Tool: Lesson Opener Template"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Before/After examples
    # Before box (red tint)
    before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(6), Inches(1.4))
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = RGBColor(255, 235, 235)
    before_box.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(5.6), Inches(1.2))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "❌ Before"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED
    p = tf.add_paragraph()
    p.text = "\"In this lesson, you'll learn about interest\""
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    p = tf.add_paragraph()
    p.text = "\"This level covers budgeting basics\""
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE

    # After box (green tint)
    after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(4.0), Inches(6), Inches(1.4))
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = RGBColor(235, 255, 235)
    after_box.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(7.033), Inches(4.1), Inches(5.6), Inches(1.2))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✅ After"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p = tf.add_paragraph()
    p.text = "\"After this, you'll know how to make your savings grow faster\""
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    p = tf.add_paragraph()
    p.text = "\"You'll be able to plan exactly where your money goes\""
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE

    # App examples
    examples_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12), Inches(1.2))
    tf = examples_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Examples from the market:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p = tf.add_paragraph()
    p.text = "• Money Wise Game: \"Life skills, not just financial skills\" | KidVestors: \"Earn while you learn\""
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_BLUE
    p = tf.add_paragraph()
    p.text = "• Blackbullion: Student money with bursary access | Bloom: \"Build wealth early\""
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_BLUE

    return slide

def add_action_2_slide(prs):
    """Action 2: Reframe to Their World (CONTEXTUAL FRAMING) - with 2x2"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar - Green for CONTEXTUAL FRAMING
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Action 2: Reframe to Their World"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Principle tag
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.4))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Based on: CONTEXTUAL FRAMING — \"Does this connect to my world?\""
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = GREEN

    # 2x2 Matrix
    # Draw grid lines
    grid_left = 1.5
    grid_top = 2.0
    grid_width = 10
    grid_height = 4.2
    mid_x = grid_left + grid_width / 2
    mid_y = grid_top + grid_height / 2

    # Vertical axis label (left)
    v_label = slide.shapes.add_textbox(Inches(0.2), Inches(mid_y - 0.5), Inches(1.2), Inches(1))
    tf = v_label.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "OUTCOME"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Horizontal axis label (bottom)
    h_label = slide.shapes.add_textbox(Inches(mid_x - 1), Inches(grid_top + grid_height + 0.1), Inches(2), Inches(0.4))
    tf = h_label.text_frame
    p = tf.paragraphs[0]
    p.text = "FRAMING"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Axis end labels
    # Top: "Action"
    top_label = slide.shapes.add_textbox(Inches(0.3), Inches(grid_top), Inches(1), Inches(0.3))
    tf = top_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Action"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE

    # Bottom: "Knowledge"
    bottom_label = slide.shapes.add_textbox(Inches(0.15), Inches(grid_top + grid_height - 0.3), Inches(1.2), Inches(0.3))
    tf = bottom_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Knowledge"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE

    # Left: "Their World"
    left_label = slide.shapes.add_textbox(Inches(grid_left), Inches(grid_top + grid_height + 0.1), Inches(1.5), Inches(0.3))
    tf = left_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Their World"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE

    # Right: "Adult World"
    right_label = slide.shapes.add_textbox(Inches(grid_left + grid_width - 1.5), Inches(grid_top + grid_height + 0.1), Inches(1.5), Inches(0.3))
    tf = right_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Adult World"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.RIGHT

    # Four quadrants
    quad_w = grid_width / 2 - 0.1
    quad_h = grid_height / 2 - 0.1

    quadrants = [
        {'name': 'EMPOWERING', 'x': grid_left, 'y': grid_top, 'color': GREEN,
         'example': '"Next time you get birthday money, try putting half in savings"'},
        {'name': 'PRESCRIPTIVE', 'x': mid_x + 0.1, 'y': grid_top, 'color': ORANGE,
         'example': '"Save 20% of your income"'},
        {'name': 'RELATABLE', 'x': grid_left, 'y': mid_y + 0.1, 'color': ACCENT_BLUE,
         'example': '"When your £10 becomes £11, that\'s interest working"'},
        {'name': 'TEXTBOOK', 'x': mid_x + 0.1, 'y': mid_y + 0.1, 'color': LIGHT_GRAY,
         'example': '"Interest is the cost of borrowing money"'},
    ]

    for q in quadrants:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(q['x']), Inches(q['y']), Inches(quad_w), Inches(quad_h))
        box.fill.solid()
        box.fill.fore_color.rgb = q['color']
        box.line.fill.background()

        # Quadrant name
        name_box = slide.shapes.add_textbox(Inches(q['x'] + 0.15), Inches(q['y'] + 0.15), Inches(quad_w - 0.3), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = q['name']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE if q['color'] != LIGHT_GRAY else DARK_BLUE

        # Example
        ex_box = slide.shapes.add_textbox(Inches(q['x'] + 0.15), Inches(q['y'] + 0.6), Inches(quad_w - 0.3), Inches(quad_h - 0.8))
        tf = ex_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = q['example']
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = WHITE if q['color'] != LIGHT_GRAY else DARK_BLUE

    # Arrow showing direction
    arrow_box = slide.shapes.add_textbox(Inches(9), Inches(2.3), Inches(2.5), Inches(0.5))
    tf = arrow_box.text_frame
    p = tf.paragraphs[0]
    p.text = "← Move here"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN

    return slide

def add_action_3_slide(prs):
    """Action 3: Elevate, Don't Simplify (CAPABILITY)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar - Purple for CAPABILITY
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Action 3: Elevate, Don't Simplify"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Principle tag
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.4))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Based on: CAPABILITY — \"Am I being talked down to?\""
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = PURPLE

    # The action
    action_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(12.333), Inches(0.9))
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = LIGHT_GRAY
    action_box.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(0.7))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Don't remove complexity — translate it. Keep the sophistication, change the language."
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Tool label
    tool_label = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(0.4))
    tf = tool_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Tool: Jargon-to-Relational Swap List"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Swap table
    swaps = [
        ('Interest', 'How your money grows while you sleep'),
        ('Compound interest', 'Growth on top of growth'),
        ('Diversification', "Don't bet everything on one thing"),
        ('Credit score', 'Your money reputation'),
        ('Budget', 'Your money plan'),
        ('Liability', 'Something that costs you money'),
        ('Asset', 'Something that makes you money'),
    ]

    # Table headers
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.5), Inches(4), Inches(0.5))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = PURPLE
    header_left.line.fill.background()

    hl_text = slide.shapes.add_textbox(Inches(0.6), Inches(3.55), Inches(3.8), Inches(0.4))
    tf = hl_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Jargon"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.6), Inches(3.5), Inches(8.233), Inches(0.5))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = PURPLE
    header_right.line.fill.background()

    hr_text = slide.shapes.add_textbox(Inches(4.7), Inches(3.55), Inches(8), Inches(0.4))
    tf = hr_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Elevated Translation"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table rows
    y = 4.05
    for i, (jargon, translation) in enumerate(swaps):
        row_color = LIGHT_GRAY if i % 2 == 0 else WHITE

        left_cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(4), Inches(0.42))
        left_cell.fill.solid()
        left_cell.fill.fore_color.rgb = row_color
        left_cell.line.fill.background()

        left_text = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.08), Inches(3.8), Inches(0.3))
        tf = left_text.text_frame
        p = tf.paragraphs[0]
        p.text = jargon
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        right_cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.6), Inches(y), Inches(8.233), Inches(0.42))
        right_cell.fill.solid()
        right_cell.fill.fore_color.rgb = row_color
        right_cell.line.fill.background()

        right_text = slide.shapes.add_textbox(Inches(4.7), Inches(y + 0.08), Inches(8), Inches(0.3))
        tf = right_text.text_frame
        p = tf.paragraphs[0]
        p.text = f'"{translation}"'
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_BLUE

        y += 0.44

    return slide

def add_action_4_slide(prs):
    """Action 4: Purpose-Driven Interactions (UTILITY) - OUT OF SCOPE"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar - Blue but greyed out
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(180, 180, 180)  # Grey
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Action 4: Purpose-Driven Interactions"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Out of scope badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.25), Inches(2.5), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(100, 100, 100)
    badge.line.fill.background()

    badge_text = slide.shapes.add_textbox(Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.5))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.text = "V2 / OUT OF SCOPE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Principle tag (greyed)
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.4))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Based on: UTILITY — \"Does this interaction have a point?\""
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 150, 150)

    # Why parked box
    why_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.1), Inches(12.333), Inches(1.8))
    why_box.fill.solid()
    why_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    why_box.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(1.6))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Why this is parked for V2:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p = tf.add_paragraph()
    p.text = "UTILITY is about interaction design — ensuring every tap, quiz, or feedback loop drives toward a functional outcome. This requires changes to mechanics and flows, not just copy."
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p = tf.add_paragraph()
    p.text = "For this first round (minor text evolutions), the three copy-focused actions are actionable."
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(100, 100, 100)

    # What it would look like
    future_label = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(0.4))
    tf = future_label.text_frame
    p = tf.paragraphs[0]
    p.text = "When this becomes relevant (V2):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)

    future_examples = [
        "Audit each interaction: \"What can the learner DO after this?\"",
        "Virtual decisions → visible consequences (Money Wise Game model)",
        "Simulations with measurable impact (Kredit Academy credit score)",
        "Complete functional chains: action → outcome → reward",
    ]

    y = 4.7
    for example in future_examples:
        ex_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(12), Inches(0.4))
        tf = ex_box.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + example
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(150, 150, 150)
        y += 0.4

    return slide

# =====================================================
# BUILD THE DECK
# =====================================================

# 1. Title slide
add_title_slide(prs, "Financial Literacy Landscape Analysis", "Modalities, Players & Strategic Directions for Money Monsters")

# 2. Executive Summary
add_content_slide(prs, "Executive Summary", [
    "Analysed 20+ solutions across the youth financial literacy space",
    "Solutions cluster into FOUR approaches: TEACH, PRACTICE, CONNECT, GUIDE",
    "Money Monsters sits in TEACH (app-based, gamified) — strong for 13-15s, weak for 18-25s",
    "Most viable path: Add GUIDE (AI coach) as unifying layer, then consider PRACTICE or CONNECT",
    "Key warning: Young adults reject 'childish gamification' — serving full age range requires evolution"
])

# 3. Section: The Landscape
add_section_header(prs, "Part 1: The Current Landscape", "What solutions exist and how they differ")

# 4. Four Approaches Framework
directions = [
    {
        'label': 'TEACH',
        'name': 'Structured Knowledge Building',
        'tagline': 'Build understanding through curriculum and content',
        'best_for': 'Young Teens (13-15)'
    },
    {
        'label': 'PRACTICE',
        'name': 'Real-World Application',
        'tagline': 'Learn by doing with real money',
        'best_for': 'Young Adults (18-25)'
    },
    {
        'label': 'CONNECT',
        'name': 'Social & Peer Learning',
        'tagline': 'Learn from and with peers',
        'best_for': 'Young & Older Teens'
    },
    {
        'label': 'GUIDE',
        'name': 'Personalized Support',
        'tagline': 'Get adaptive, personalized support',
        'best_for': 'All ages, especially 16+'
    }
]
add_framework_slide(prs, "Four Approaches to Youth Financial Literacy", directions)

# 5-8. Approach detail slides
direction_details = [
    {
        'label': 'TEACH',
        'name': 'Structured Knowledge Building',
        'tagline': 'Build understanding through curriculum and content',
        'best_for': 'Young Teens (13-15)',
        'mechanics': [
            'Progressive curriculum with clear learning path',
            'Gamification: streaks, points, badges, levels',
            'Quizzes and knowledge checks',
            'Character-driven narratives (e.g. monsters)',
            'School/institutional delivery options'
        ],
        'examples': [
            'Money Monsters, Zogo (app-based)',
            'Money Wise Game, Stock Market Game (simulation)',
            'Money Ready, Young Enterprise (institutional)',
            'Kredit Academy, KidVestors',
            'MoneyTime Kids, Blackbullion'
        ]
    },
    {
        'label': 'PRACTICE',
        'name': 'Real-World Application',
        'tagline': 'Learn by doing with real money',
        'best_for': 'Young Adults (18-25)',
        'mechanics': [
            'Banking app with educational layer',
            'Real money management with oversight',
            'Behavioural nudges at spending moments',
            'Automatic saving and round-up features',
            'Investment with fractional shares'
        ],
        'examples': [
            'Greenlight (6M+ users)',
            'GoHenry Money Missions',
            'Gimi, Mega Card, Earlybean',
            'Bloom, Stockpile (teen investing)',
            'Acorns round-ups'
        ]
    },
    {
        'label': 'CONNECT',
        'name': 'Social & Peer Learning',
        'tagline': 'Learn from and with peers',
        'best_for': 'Young & Older Teens',
        'mechanics': [
            'Challenge friends / social competition',
            'Community discussion and forums',
            'Peer mentoring (near-peer model)',
            'Finfluencer content and social proof',
            'Ambassador / peer teaching programs'
        ],
        'examples': [
            'TikTok finfluencers (Vivian Tu, etc.)',
            'Bumper "Investment Circle"',
            'Discord/Reddit finance communities',
            'Futures Financials (teen-led)',
            'Finance podcasts for youth'
        ]
    },
    {
        'label': 'GUIDE',
        'name': 'Personalized Support',
        'tagline': 'Get adaptive, personalized support',
        'best_for': 'All ages, especially 16+',
        'mechanics': [
            'Conversational AI (ask anything)',
            'Adaptive difficulty based on knowledge',
            'Personalised learning paths',
            '24/7 availability for just-in-time help',
            'Proactive nudges based on patterns'
        ],
        'examples': [
            'Cleo (AI spending assistant)',
            'Quirk (personality-based advice)',
            'Coach mAIa, Kiro AI Coach',
            'Tendi AI financial advisor',
            'ChatGPT for finance questions'
        ]
    }
]

colors = [ORANGE, ACCENT_BLUE, GREEN, PURPLE]
for direction, color in zip(direction_details, colors):
    add_direction_detail_slide(prs, direction, color)

# 9. Section: Approach Headers
add_section_header(prs, "Part 2: The Four Approaches", "TEACH / PRACTICE / CONNECT / GUIDE")

# 10. Modality headers overview
add_modality_headers_slide(prs)

# 11-14. Apps by approach
teach_apps = [
    {'name': 'Money Monsters', 'desc': 'Gamified app with monster-taming metaphor (app-based)'},
    {'name': 'Zogo', 'desc': 'Gamified financial literacy with real rewards (app-based)'},
    {'name': 'Money Wise Game / Dot Dot Fire', 'desc': 'Educational games and simulations'},
    {'name': 'Kredit Academy / KidVestors', 'desc': 'Gamified courses with progression'},
    {'name': 'Money Ready / Young Enterprise', 'desc': 'School-delivered curriculum (institutional)'},
    {'name': 'Blackbullion / MoneyTime', 'desc': 'Online courses for students (institutional)'},
]
add_apps_by_modality_slide(prs, "TEACH: Structured Knowledge Building", teach_apps, ORANGE)

practice_apps = [
    {'name': 'Greenlight', 'desc': '6M+ users, debit card + savings + gamified learning'},
    {'name': 'GoHenry', 'desc': 'Money Missions with real money management'},
    {'name': 'Gimi / Mega Card', 'desc': 'Banking as learning tool, parental oversight'},
    {'name': 'Bloom / Stockpile', 'desc': 'Teen investing with fractional shares'},
    {'name': 'Acorns', 'desc': 'Round-up investing, automatic saving'},
    {'name': 'Earlybean / BusyKid', 'desc': 'Pocket money apps with real money'},
]
add_apps_by_modality_slide(prs, "PRACTICE: Real-World Application", practice_apps, PURPLE)

connect_apps = [
    {'name': 'TikTok Finfluencers', 'desc': 'Vivian Tu, Taylor Price, Kia Commodore'},
    {'name': 'Bumper Investment Circle', 'desc': 'Social learning with peer accountability'},
    {'name': 'Discord/Reddit Communities', 'desc': 'r/personalfinance, finance servers'},
    {'name': 'Futures Financials', 'desc': 'Teen-led peer teaching programs'},
    {'name': 'Finance Podcasts', 'desc': 'Money Moves, Teen Money Matters, etc.'},
]
add_apps_by_modality_slide(prs, "CONNECT: Social & Peer Learning", connect_apps, GREEN)

guide_apps = [
    {'name': 'Cleo', 'desc': 'AI spending assistant with personality'},
    {'name': 'Quirk', 'desc': 'Personality-based financial advice'},
    {'name': 'Coach mAIa / Kiro', 'desc': 'AI coaching with personalized budgeting'},
    {'name': 'Tendi', 'desc': 'AI financial advisor for beginners'},
    {'name': 'ChatGPT / Claude', 'desc': 'General AI for finance questions'},
]
add_apps_by_modality_slide(prs, "GUIDE: Personalized Support", guide_apps, ACCENT_BLUE)

# 15. Section: Engagement Principles
add_section_header(prs, "Part 3: Engagement Principles", "How to deliver financial literacy effectively")

# 16. Engagement principles slide
add_engagement_principles_slide(prs)

# 17-20. Action slides
add_action_1_slide(prs)  # Lead with Outcomes (PROPOSITION)
add_action_2_slide(prs)  # Reframe to Their World (CONTEXTUAL FRAMING)
add_action_3_slide(prs)  # Elevate, Don't Simplify (CAPABILITY)
add_action_4_slide(prs)  # Purpose-Driven Interactions (UTILITY) - greyed out

# 21. Section: Audience Fit
add_section_header(prs, "Part 4: Audience Fit Analysis", "Which approach works for which age group?")

# 22. Audience fit matrix
add_audience_fit_slide(prs)

# 23. Section: So What
add_section_header(prs, "Part 5: The 'So What'", "Implications for Money Monsters")

# 24. So What slide
add_so_what_slide(prs)

# 25. Final recommendations
add_content_slide(prs, "Recommended Actions", [
    "Short-term: Partner with finfluencers for authentic reach (CONNECT elements)",
    "Short-term: Create TikTok-ready shareable content from app",
    "Medium-term: Add 'ask anything' AI chat feature (GUIDE)",
    "Medium-term: Add friend challenges and social features (CONNECT)",
    "Longer-term: Explore banking partnerships (PRACTICE)",
    "Premium opportunity: GUIDE features, CONNECT community, PRACTICE integration"
])

# Save
output_path = '/Users/mattspeak/Library/Mobile Documents/iCloud~md~obsidian/Documents/ObsidianMatt/WORK/Mosaic/Money Monsters User Testing Feb-26/Research/Desk_research/Financial_Literacy_Landscape_Deck.pptx'
prs.save(output_path)
print(f"Deck saved to: {output_path}")
